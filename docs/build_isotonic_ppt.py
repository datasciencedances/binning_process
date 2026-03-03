"""
Script tạo file PowerPoint (.pptx) cho slide Isotonic Regression Binning.
Chạy: python docs/build_isotonic_ppt.py
Cần: pip install python-pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

# Đường dẫn output
OUTPUT_PATH = Path(__file__).resolve().parent / "Isotonic_Regression_Binning.pptx"


def add_title_slide(prs, title: str, subtitle: str = ""):
    """Slide layout 0 = title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title: str, bullets: list):
    """Slide layout 1 = title and content."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(bullets):
        p = body.add_paragraph() if i == 0 else body.add_paragraph()
        p.text = item.replace("**", "").strip()
        p.level = 0
        p.space_after = Pt(6)
    return slide


def add_slide_with_table(prs, title: str, headers: list, rows: list):
    """Slide có tiêu đề + bảng."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    left = Inches(0.5)
    top = Inches(0.4)
    w, h = Inches(9), Inches(0.8)
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    # Table
    cols, row_count = len(headers), 1 + len(rows)
    tw = 9.0 / cols
    table_top = Inches(1.3)
    table = slide.shapes.add_table(row_count, cols, Inches(0.5), table_top, Inches(9), Inches(0.4 * row_count)).table
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        table.cell(0, c).text_frame.paragraphs[0].font.bold = True
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            if c < len(row):
                table.cell(r + 1, c).text = cell_text
    return slide


def add_slide_with_code(prs, title: str, code_lines: list, bullets: list):
    """Slide có tiêu đề + khối code (text) + bullet."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left, top = Inches(0.5), Inches(0.4)
    tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.7))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(28)
    tx.text_frame.paragraphs[0].font.bold = True
    # Code block
    code_top = Inches(1.2)
    code_box = slide.shapes.add_textbox(Inches(0.5), code_top, Inches(9), Inches(1.8))
    tf = code_box.text_frame
    tf.word_wrap = True
    for line in code_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(11)
    # Bullets
    bullet_top = Inches(3.2)
    bullet_box = slide.shapes.add_textbox(Inches(0.5), bullet_top, Inches(9), Inches(3.5))
    bf = bullet_box.text_frame
    for i, item in enumerate(bullets):
        p = bf.add_paragraph()
        p.text = item.replace("**", "").strip()
        p.space_after = Pt(4)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Tiêu đề ---
    add_title_slide(
        prs,
        "Isotonic Regression Binning",
        "Binning đơn điệu dựa trên hồi quy đẳng áp\n\n"
        "• Ứng dụng: Scorecard, risk modeling, credit scoring\n"
        "• Đảm bảo: Biến bin có quan hệ đơn điệu với event rate\n"
        "• Công cụ: Isotonic Regression (sklearn) + quantile chia sơ bộ",
    )

    # --- Slide 2: Binning là gì ---
    add_content_slide(
        prs,
        "Binning trong modeling",
        [
            "Binning: Chia biến liên tục X thành các khoảng (bins), mỗi khoảng gán một giá trị (index, WOE, score).",
            "Mục đích: Giảm nhiễu, dễ giải thích; xử lý phi tuyến (X ↔ risk); chuẩn hóa cho scorecard (WOE, điểm).",
            "Thách thức: Chọn cut-points sao cho vừa đủ bins, vừa ổn định, vừa có ý nghĩa (ví dụ đơn điệu).",
        ],
    )

    # --- Slide 3: Monotonic ---
    add_content_slide(
        prs,
        "Tại sao cần monotonic?",
        [
            "Monotonic: Khi X tăng (hoặc giảm), event rate chỉ đi một chiều — tăng dần hoặc giảm dần.",
            "Lý do: Giải thích dễ với business/compliance; ổn định (tránh bin lồi lõm do nhiễu); regulation (fair lending).",
            "Isotonic Binner: Tạo cut-points sao cho sau khi bin, event rate theo bin luôn đơn điệu.",
        ],
    )

    # --- Slide 4: Isotonic Regression là gì ---
    add_content_slide(
        prs,
        "Isotonic Regression (hồi quy đẳng áp)",
        [
            "Định nghĩa: Tìm hàm f(x) đơn điệu (tăng hoặc giảm) xấp xỉ dữ liệu tốt nhất (tổng bình phương sai số có trọng số).",
            "Ràng buộc: Chỉ đi lên (increasing) hoặc đi xuống (decreasing), không lên xuống tùy ý.",
            "Hình ảnh: Đường bậc thang — bậc ngang = mức giá trị bằng nhau; chỗ gãy = ranh giới (cut-point).",
            "Trong binning: Input = (center bin, event rate); output = đường bậc thang đơn điệu; điểm gãy → cut-point trên X.",
        ],
    )

    # --- Slide 5: Ý tưởng trực quan ---
    add_slide_with_code(
        prs,
        "Từ dữ liệu nhiễu → đường đơn điệu",
        [
            "Event rate",
            "    ^",
            "    |     *  *",
            "    |   *   *    *   ← Dữ liệu thô lên xuống do nhiễu",
            "    | *       *     *",
            "    +-------------------------> X (bin center)",
            "",
            "Isotonic vẽ đường bậc thang đi LÊN hoặc XUỐNG qua các điểm.",
        ],
        [
            "Mỗi bậc ngang = một nhóm bin có cùng event rate (đã làm mượt).",
            "Điểm gãy giữa hai bậc = cut-point trên trục X → chia bin cuối cùng.",
        ],
    )

    # --- Slide 6: Thuật toán 4 bước (bảng) ---
    add_slide_with_table(
        prs,
        "Các bước của Isotonic Binner",
        ["Bước", "Mô tả"],
        [
            ("1", "Chia X theo quantile (vd. 20 bins) → init_cuts, edges, bin_idx."),
            ("2", "Mỗi bin: center (median X), event rate, weight (số lượng)."),
            ("3", "Fit Isotonic Regression: (centers, rates, weights) → đường bậc thang đơn điệu."),
            ("4", "Điểm gãy: hai bậc liên tiếp khác nhau → boundary → giá trị X thực → cut. Giới hạn max_bins."),
        ],
    )

    # --- Slide 7: Bước 1 & 2 ---
    add_content_slide(
        prs,
        "Chi tiết bước 1 và 2",
        [
            "Bước 1 — Quantile: quantile_cuts(x, n_init_bins) → edges = [-∞, cuts..., +∞], bin_idx = pd.cut(x, edges).",
            "Bước 2 — Đại diện mỗi bin: Center = median(x[mask]); Rate = y[mask].mean(); Weight = mask.sum().",
            "→ Bộ ba (centers, rates, weights) làm đầu vào cho Isotonic Regression.",
        ],
    )

    # --- Slide 8: Bước 3 & 4 ---
    add_content_slide(
        prs,
        "Chi tiết bước 3 và 4",
        [
            "Bước 3 — Fit: IsotonicRegression(increasing=(direction_=='ascending')); fitted = iso.fit_transform(centers, rates, sample_weight=weights).",
            "Bước 4 — Cut-points: Duyệt (fitted[i], fitted[i+1]); nếu khác nhau → boundary = (centers[i]+centers[i+1])/2 → ánh xạ về X thực → cuts. Cắt bớt theo max_bins.",
        ],
    )

    # --- Slide 9: Ưu nhược ---
    add_content_slide(
        prs,
        "Ưu điểm & Hạn chế",
        [
            "Ưu điểm: Đơn điệu đảm bảo; tận dụng center/rate/weight có trọng số; ổn định (làm mượt nhiễu); dùng sklearn, dễ bảo trì.",
            "Hạn chế: Phụ thuộc n_init_bins (quá ít → ít điểm gãy; quá nhiều → overfit); chỉ mô hình một chiều, không dạng U.",
        ],
    )

    # --- Slide 10: Tóm tắt ---
    add_content_slide(
        prs,
        "Tóm tắt",
        [
            "Isotonic Binner dùng Isotonic Regression để tìm đường bậc thang đơn điệu qua (center, event rate) của các bin quantile.",
            "Cut-points = các điểm gãy của đường bậc thang, ánh xạ về giá trị X.",
            "Đảm bảo monotonic sau khi bin, phù hợp scorecard và quy định.",
            "Luồng: Quantile → centers/rates/weights → Isotonic fit → điểm gãy → cuts (max_bins).",
            "Code: binning_process/supervised/isotoic.py — class IsotonicBinner.",
        ],
    )

    prs.save(OUTPUT_PATH)
    print("Saved:", str(OUTPUT_PATH))


if __name__ == "__main__":
    main()
